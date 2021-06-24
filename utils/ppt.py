from pptx import Presentation

# Todo: make lyrics version in normal font
# Todo: Add structure text to top left placeholder for slides that have been split in two

# 0 = new slide indicator
# 1 = paragraph type (i.e., chorus, verse)
# 2 = chords
# 3 = lyrics
# 4 = paragraph type for chords only
# 5 = lyrics line for chord ppt only
# 6 = lyrics line for lyrics ppt only
# 7 = new slide indicator only chords ppt
# 8 = new slide indicator only lyrics ppt


def create_lyrics_ppt(data):
    """
    Creates the lyrics ppt file. This excludes slides for chords only
    structure elements (e.g., Instrumental, Interlude) and excludes
    the chords lines.
    :param data: The complete data for all selected songs
    """
    my_presentation = Presentation("dewg_template_lyrics.pptx")
    for song in data:
        slides = song[0]
        title = song[1]
        # artist = song[2]
        # key = song[3]

        # initialize variables
        no_reset_flag = False
        skip_flag = False
        content_string = ""
        structure_text = ""

        for slide in slides:
            if no_reset_flag:
                no_reset_flag = False  # prevents the content and structure strings to be reset if True
            else:
                content_string = ""
                structure_text = ""

            for line in slide:
                if line[0] == 4:
                    # skip slide of this type (i.e., interlude, instrumental)
                    pass
                elif line[0] == 1:
                    structure_text += line[1]
                    structure_text += " & "
                    pass
                elif line[0] == 2 or line[0] == 5:
                    # skip chord lines
                    # skip lyrics line for chord only
                    pass
                elif line[0] == 7:
                    # skip new slide indicator for chords ppt
                    # do not advance to next slide, but combine current and next slide into one
                    no_reset_flag = True  # stops new slides from being added and stops strings resetting
                else:
                    content_string += line[1]
                    content_string += "\n"

            if not no_reset_flag:
                # create new slide
                this_slide = my_presentation.slides.add_slide(my_presentation.slide_layouts[0])
                # add song title to top right placeholder
                this_slide.placeholders[10].text = title
                # add main content to the slide
                if content_string.startswith('\n'):
                    this_slide.placeholders[1].text = content_string[1:].rstrip()
                else:
                    this_slide.placeholders[1].text = content_string.rstrip()
                this_slide.placeholders[0].text = structure_text[:-3]

    my_presentation.save("presentation_lyrics.pptx")


def create_chords_ppt(data):
    """
    Creates the chords ppt file. This includes slides for chords only
    structure elements (e.g., Instrumental, Interlude) and includes
    the chords lines.
    :param data: The complete data for all selected songs
    """

    # Todo: Add current song key to top right corner
    # Todo: Make lyrics text gray to focus on chords ->
    #  https://stackoverflow.com/questions/24729098/text-color-in-python-pptx-module

    my_presentation = Presentation("dewg_template_chords.pptx")
    for song in data:
        slides = song[0]
        title = song[1]
        # artist = song[2]
        # key = song[3]

        # initialize variables
        no_reset_flag = False
        content_string = ""
        structure_text = ""

        for slide in slides:
            if no_reset_flag:
                no_reset_flag = False  # prevents the content and structure strings to be reset if True
            else:
                content_string = ""
                structure_text = ""

            for line in slide:
                if line[0] == 6:
                    # skip line if it is lyrics for lyrics ppt only
                    pass
                elif line[0] == 8:
                    # skip new slide indicator for chords ppt
                    # do not advance to next slide, but combine current and next slide into one
                    no_reset_flag = True  # stops new slides from being added and stops strings resetting
                elif line[0] == 1 or line[0] == 4:
                    structure_text += line[1]
                    structure_text += " & "
                else:
                    content_string += line[1]
                    content_string += "\n"

            if not no_reset_flag:
                # create new slide
                this_slide = my_presentation.slides.add_slide(my_presentation.slide_layouts[0])
                # add song title to top right placeholder
                this_slide.placeholders[10].text = title
                # add main content to the slide
                if content_string.startswith('\n'):
                    this_slide.placeholders[1].text = content_string[1:].rstrip()
                else:
                    this_slide.placeholders[1].text = content_string.rstrip()
                this_slide.placeholders[0].text = structure_text[:-3]

    my_presentation.save("presentation_chords.pptx")
